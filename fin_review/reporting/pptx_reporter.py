"""PowerPoint presentation generation module."""

import pandas as pd
import structlog
from pathlib import Path
from typing import Dict, Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import datetime

logger = structlog.get_logger()


class PowerPointReporter:
    """Generates PowerPoint executive presentation."""
    
    def __init__(self, output_path: Path, config: Optional[Dict] = None):
        """
        Initialize PowerPoint reporter.
        
        Args:
            output_path: Path to output PPTX file
            config: Configuration dictionary
        """
        self.output_path = output_path
        self.config = config or {}
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def generate_report(
        self,
        commentary: Dict,
        kpis: Dict,
        trends: Dict,
        aging: Dict,
        anomalies: Dict
    ):
        """
        Generate complete PowerPoint presentation.
        
        Args:
            commentary: Commentary results
            kpis: KPI results
            trends: Trend results
            aging: Aging results
            anomalies: Anomaly results
        """
        logger.info(f"Generating PowerPoint report: {self.output_path}")
        
        # Title slide
        self._create_title_slide()
        
        # Executive summary
        self._create_executive_summary_slide(commentary)
        
        # Key insights
        self._create_insights_slide(commentary)
        
        # Financial overview
        self._create_financial_overview_slide(kpis)
        
        # Trends
        self._create_trends_slide(kpis)
        
        # Aging analysis
        self._create_aging_slide(aging)
        
        # Risks
        self._create_risks_slide(commentary)
        
        # Recommendations
        self._create_recommendations_slide(commentary)
        
        # Save presentation
        self.prs.save(self.output_path)
        
        logger.info(f"PowerPoint report generated successfully: {self.output_path}")
    
    def _create_title_slide(self):
        """Create title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = "Financial Review\nAnalytical Report"
        
        title_para = text_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %d, %Y")
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(20)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _create_executive_summary_slide(self, commentary: Dict):
        """Create executive summary slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and content
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Add summary text
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Get summary
        exec_summary = commentary.get('executive_summary', '')
        
        # Split into sections
        for line in exec_summary.split('\n'):
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.level = 1 if line.isupper() else 0
        
        # Add speaker notes if enabled
        if self.config.get('include_speaker_notes', True):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = exec_summary
    
    def _create_insights_slide(self, commentary: Dict):
        """Create key insights slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Key Insights"
        
        # Add insights
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        insights = commentary.get('insights', [])
        
        for i, insight in enumerate(insights[:5], 1):
            # Title
            p = text_frame.add_paragraph()
            confidence = f" [{insight['confidence'].upper()}]" if self.config.get('confidence_levels', True) else ""
            p.text = f"{i}. {insight['title']}{confidence}"
            p.font.bold = True
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            
            # Content
            p = text_frame.add_paragraph()
            p.text = insight['content']
            p.font.size = Pt(11)
            p.level = 1
            p.space_after = Pt(12)
    
    def _create_financial_overview_slide(self, kpis: Dict):
        """Create financial overview slide with key metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Financial Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        summary_kpis = kpis.get('summary_kpis', {})
        growth_metrics = kpis.get('growth_metrics', {})
        
        # Create metric boxes
        metrics = [
            ("Total Revenue", summary_kpis.get('total_revenue', 0), "€"),
            ("Total OPEX", summary_kpis.get('total_opex', 0), "€"),
            ("Net Profit", summary_kpis.get('net_profit', 0), "€"),
            ("YoY Growth", growth_metrics.get('latest_revenue_yoy', 0), "%")
        ]
        
        box_width = Inches(2)
        box_height = Inches(1.5)
        spacing = Inches(0.3)
        start_left = Inches(0.5)
        top = Inches(1.5)
        
        for i, (label, value, unit) in enumerate(metrics):
            left = start_left + (box_width + spacing) * i
            
            # Background box
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._get_color(i)
            
            # Label
            label_box = slide.shapes.add_textbox(left, top + Inches(0.2), box_width, Inches(0.4))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.alignment = PP_ALIGN.CENTER
            label_para.font.size = Pt(12)
            label_para.font.bold = True
            
            # Value
            value_box = slide.shapes.add_textbox(left, top + Inches(0.7), box_width, Inches(0.6))
            value_frame = value_box.text_frame
            
            if unit == "€":
                value_frame.text = f"€{value/1000:.1f}K"
            else:
                value_frame.text = f"{value:.1f}%"
            
            value_para = value_frame.paragraphs[0]
            value_para.alignment = PP_ALIGN.CENTER
            value_para.font.size = Pt(20)
            value_para.font.bold = True
    
    def _create_trends_slide(self, kpis: Dict):
        """Create trends slide with chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Revenue & OPEX Trends"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        monthly_kpis = kpis.get('monthly_kpis')
        
        # Convert list to DataFrame if needed
        if isinstance(monthly_kpis, list):
            if len(monthly_kpis) == 0:
                return
            monthly_kpis = pd.DataFrame(monthly_kpis)
        
        if monthly_kpis is not None and len(monthly_kpis) > 0:
            # Create chart data
            chart_data = CategoryChartData()
            
            # Handle datetime formatting
            if pd.api.types.is_datetime64_any_dtype(monthly_kpis['year_month']):
                chart_data.categories = [d.strftime('%Y-%m') for d in monthly_kpis['year_month']]
            else:
                chart_data.categories = [str(d) for d in monthly_kpis['year_month']]
            
            if 'revenue' in monthly_kpis.columns:
                chart_data.add_series('Revenue', monthly_kpis['revenue'].tolist())
            
            if 'opex' in monthly_kpis.columns:
                chart_data.add_series('OPEX', monthly_kpis['opex'].tolist())
            
            # Add chart
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            
            chart.has_legend = True
            chart.legend.position = 2  # Right
    
    def _create_aging_slide(self, aging: Dict):
        """Create aging analysis slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Receivables & Payables Aging"
        
        # Add aging summary
        ar_summary = aging.get('ar_summary', {})
        ap_summary = aging.get('ap_summary', {})
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        
        # AR Summary
        p = text_frame.add_paragraph()
        p.text = "Accounts Receivable"
        p.font.bold = True
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        ar_outstanding = ar_summary.get('total_outstanding', 0)
        ar_overdue_pct = ar_summary.get('overdue_pct', 0)
        
        p = text_frame.add_paragraph()
        p.text = f"• Total Outstanding: €{ar_outstanding/1000:.1f}K"
        p.font.size = Pt(12)
        
        p = text_frame.add_paragraph()
        p.text = f"• Overdue: {ar_overdue_pct:.1f}%"
        p.font.size = Pt(12)
        p.space_after = Pt(20)
        
        # AP Summary
        p = text_frame.add_paragraph()
        p.text = "Accounts Payable"
        p.font.bold = True
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        ap_outstanding = ap_summary.get('total_outstanding', 0)
        ap_overdue_pct = ap_summary.get('overdue_pct', 0)
        
        p = text_frame.add_paragraph()
        p.text = f"• Total Outstanding: €{ap_outstanding/1000:.1f}K"
        p.font.size = Pt(12)
        
        p = text_frame.add_paragraph()
        p.text = f"• Overdue: {ap_overdue_pct:.1f}%"
        p.font.size = Pt(12)
    
    def _create_risks_slide(self, commentary: Dict):
        """Create risks slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Top Risks"
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        risks = commentary.get('risks', [])
        
        for i, risk in enumerate(risks[:5], 1):
            # Title
            p = text_frame.add_paragraph()
            confidence = f" [{risk['confidence'].upper()}]" if self.config.get('confidence_levels', True) else ""
            p.text = f"{i}. {risk['title']}{confidence}"
            p.font.bold = True
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            
            # Content
            p = text_frame.add_paragraph()
            p.text = risk['content']
            p.font.size = Pt(11)
            p.level = 1
            p.space_after = Pt(12)
    
    def _create_recommendations_slide(self, commentary: Dict):
        """Create recommendations slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Recommended Actions"
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        recommendations = commentary.get('recommendations', [])
        
        for i, rec in enumerate(recommendations[:5], 1):
            # Title
            p = text_frame.add_paragraph()
            p.text = f"{i}. {rec['title']}"
            p.font.bold = True
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            
            # Content
            p = text_frame.add_paragraph()
            p.text = rec['content']
            p.font.size = Pt(11)
            p.level = 1
            p.space_after = Pt(12)
    
    @staticmethod
    def _get_color(index: int):
        """Get color for metric box."""
        colors = [
            RGBColor(79, 129, 189),   # Blue
            RGBColor(155, 187, 89),   # Green
            RGBColor(192, 80, 77),    # Red
            RGBColor(128, 100, 162),  # Purple
        ]
        return colors[index % len(colors)]


def generate_pptx_report(
    output_path: Path,
    commentary: Dict,
    kpis: Dict,
    trends: Dict,
    aging: Dict,
    anomalies: Dict,
    config: Optional[Dict] = None
):
    """
    Convenience function to generate PowerPoint report.
    
    Args:
        output_path: Path to output PPTX file
        commentary: Commentary results
        kpis: KPI results
        trends: Trend results
        aging: Aging results
        anomalies: Anomaly results
        config: Configuration dictionary
    """
    reporter = PowerPointReporter(output_path, config)
    reporter.generate_report(commentary, kpis, trends, aging, anomalies)

